import os
from pathlib import Path
from textwrap import wrap

try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception:
    print("Missing python-pptx. Install with: pip install python-pptx")
    raise


ROOT = Path(__file__).resolve().parents[1]


def read_text(path):
    try:
        return Path(path).read_text(encoding="utf-8")
    except Exception:
        return ""


def short_lines(text, maxlen=80, maxlines=8):
    lines = []
    for para in text.splitlines():
        para = para.strip()
        if not para:
            continue
        for w in wrap(para, maxlen):
            lines.append(w)
            if len(lines) >= maxlines:
                return lines
    return lines


def add_bullets(slide, title, bullets):
    title_box = slide.shapes.title
    title_box.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
        p.level = 0


def find_sample_image():
    candidates = [
        ROOT / "static" / "test_samples",
        ROOT / "Dataset" / "test",
        ROOT / "static" / "uploads",
    ]
    for c in candidates:
        if c.exists() and c.is_dir():
            for p in c.iterdir():
                if p.suffix.lower() in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
                    return p
    return None


def build_presentation(readme, project_expl):
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Brain Stroke Detection — AI-Based Medical Image Analysis"
    slide.placeholders[1].text = "CNN + ML models — Project Overview"

    # Abstract (from README first paragraph)
    abstract_lines = []
    for line in (readme or project_expl).splitlines():
        line = line.strip()
        if line:
            abstract_lines.append(line)
            break
    if abstract_lines:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        add_bullets(slide, "Abstract", short_lines('\n'.join(abstract_lines), maxlen=100, maxlines=4))

    # Introduction (brief)
    intro = short_lines(project_expl, maxlen=90, maxlines=6)
    if intro:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        add_bullets(slide, "Introduction", intro)

    # Problem Statement
    problem = [
        "Rapid detection of brain stroke from CT scans is critical to save brain tissue",
        "Shortage of specialists and time pressure in emergency settings",
        "Privacy and availability issues with real patient CT datasets (uses synthetic data)"
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Problem Statement", problem)

    # Existing Methodology
    existing = [
        "Manual radiologist evaluation of CT scans (time-consuming)",
        "Classical ML approaches: Random Forest, SVM, Logistic Regression (feature-based)",
        "Limited spatial understanding compared to deep learning"
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Existing Methodology", existing)

    # Proposed Methodology
    proposed = [
        "End-to-end Convolutional Neural Network trained on CT images",
        "Preprocessing: grayscale, resize to 128x128, normalization",
        "4 conv layers + dense head with dropout; outputs stroke probability",
        "Integrated into a Flask web app for fast inference and history tracking"
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Proposed Methodology", proposed)

    # Summary
    summary = short_lines(readme, maxlen=90, maxlines=6)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Project Summary", summary)

    # Features
    features = [
        "CNN-based detection (Normal / Stroke)",
        "Image upload + real-time prediction + confidence",
        "Model comparison: CNN, Random Forest, SVM, Logistic Regression",
        "Dashboard, history, authentication, Docker support",
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Key Features", features)

    # Tech stack
    tech = ["Backend: Flask, SQLite", "Deep Learning: PyTorch", "ML: scikit-learn", "Frontend: Bootstrap, Chart.js"]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Tech Stack", tech)

    # Dataset
    dataset = ["Synthetic CT images (800 train, 200 test)", "6 sample images (3 normal + 3 stroke)", "128x128 grayscale inputs"]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Dataset", dataset)

    # Architecture
    arch = [
        "4 Conv layers: 32→64→128→256 + MaxPool",
        "Flatten → Dense(256) + Dropout(0.5) → Dense(1) + Sigmoid",
        "Input: 1 x 128 x 128 (grayscale CT)"
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "CNN Architecture", arch)

    # Training & results
    training = ["Training: 15 epochs, Adam optimizer", "Results: CNN ~100% accuracy (synthetic data)", "Compared: RF 99%, SVM 97%, LR 95%"]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Training & Results", training)

    # Web app pages
    pages = ["Login/Register", "Home (stats + samples)", "Predict (upload + inference)", "History (past scans)", "Dashboard (charts)"]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Web Application", pages)

    # Quick start
    steps = ["pip install -r requirements.txt", "python generate_dataset.py", "python train_model.py", "python app.py (open http://localhost:5010)"]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Quick Start", steps)

    # Disclaimer
    disclaimer = ["Educational project with synthetic data — not for clinical use", "Consult medical professionals for diagnosis"]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Disclaimer", disclaimer)

    # Sample image
    sample = find_sample_image()
    if sample:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Sample CT Image"
            slide.shapes.add_picture(str(sample), Inches(1), Inches(1.5), width=Inches(6))
        except Exception:
            pass

    out_path = ROOT / "presentation.pptx"
    prs.save(str(out_path))
    print(f"Presentation saved to: {out_path}")


def main():
    readme = read_text(ROOT / "README.md")
    proj = read_text(ROOT / "PROJECT_EXPLANATION.md")
    build_presentation(readme or proj, proj)


if __name__ == "__main__":
    main()
